import os
import sys
import PythonMagick
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfFileWriter, PdfFileReader
from logger import Logger


class PdfToPpt(object):
    def __init__(self, pdf_file=None, ppt_file=None):
        self.pdf_file = pdf_file
        self.ppt_file = pdf_file.replace('.pdf', '.pptx')
	self.total_pages = 1
        self.log = Logger.defaults('PdfToPptx')
        self.log.debug('%s \n %s' % (self.pdf_file, self.ppt_file))

    def check_file_exist(self, file_path):
        self.log.info('Checking file - %s ' % file_path)
        if os.path.isfile(file_path):
            return True
        else:
            return False

    def pdf_to_image(self, pdf_file):
        if not self.check_file_exist(pdf_file):
            self.log.debug('Requested file not found in %s ' % pdf_file)
            return False
        image_file = pdf_file.replace('.pdf', '.jpg')
        try:
            pdf_to_img = PythonMagick.Image()
            pdf_to_img.density('200')
            pdf_to_img.read(pdf_file)
            pdf_to_img.write(image_file)
            self.log.info('Image convert passed - %s ' % image_file)
            return True
        except Exception:
            self.log.debug('Image convert failed - %s ' % image_file)
            self.log.error('', exc_info=True)
            return False

    def pdf_splitter(self):
        self.log.info('Called pdf_splitter')
        input_pdf = PdfFileReader(file(self.pdf_file, 'rb'))
        self.total_pages = input_pdf.numPages

        for page_number in range(self.total_pages):
            output = PdfFileWriter()
            output.addPage(input_pdf.getPage(page_number))
            # new filename
            new_pdf = '_%s%s' % (str(page_number+1), '.pdf')
            new_pdf = self.pdf_file.replace('.pdf', new_pdf)
            file_stream = file(new_pdf, 'wb')
            output.write(file_stream)
            file_stream.close()

            # calling pdf to image conversion
            self.pdf_to_image(new_pdf)

    def create_ppt(self):
        self.log.info('Called create_ppt')
        prs = Presentation()
        try:
            for slide_number in range(self.total_pages):
                img_path = self.pdf_file.replace('.pdf', '_%s%s' % (str(slide_number+1), '.jpg'))
                self.log.debug('%s' % img_path)
                new_slide = prs.slide_layouts[0]
                slide = prs.slides.add_slide(new_slide)
                subtitle = slide.placeholders[1]
                title = slide.shapes.title
                title.text = "Image %s " % str(slide_number+1)
                left = top = Inches(0.1)
                height = Inches(7.5)
                pic = slide.shapes.add_picture(img_path, left, top, height=height)
                prs.save(self.ppt_file)
        except IOError:
            self.log.error('error creating ppt', exc_info=True)

    def execute(self):
        self.log.info('Calling the main execution for ppt conversion')
        self.pdf_splitter()
        self.create_ppt()
        self.log.info('Done ppt conversion')


if __name__ == '__main__':
     run_time = sys.argv[1:]
     PdfToPpt(pdf_file=run_time[0]).execute()

