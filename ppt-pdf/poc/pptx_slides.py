#!/usr/bin/env python


from pptx import Presentation

prs = Presentation()

# 1st slide
title_slide_layouts = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layouts)

title = slide.shapes.title
subtitle =  slide.placeholders[1]

title.text = "Hello World!"
subtitle.text = "python-pptx was here!"


# 2nd slide
bullet_slide_layouts = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layouts)
shapes = slide.shapes

title_shape = shapes.title
body_shape =  shapes.placeholders[1]

title_shape.text = "Adding a bullet slide "

tf = body_shape.text_frame
tf.text = "Find the bullet slide layout"

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

# 3rd slide
from pptx.util import Inches

img_path = 'a.jpg'
blank_slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(0.1)
height = Inches(7.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)


prs.save('test.pptx')






